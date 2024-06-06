from langchain_community.document_loaders import PyMuPDFLoader
import docx
from pptx import Presentation

def extract_text(path):
    return open(path, 'r').read()

def extract_pdf(path):
    loader = PyMuPDFLoader(path)
    data = loader.load()
    data = [x.page_content for x in data]
    content = '\n\n'.join(data)
    return content

def extract_docx(path):
    doc = docx.Document(path)
    data = []
    for paragraph in doc.paragraphs:
        data.append(paragraph.text)
    content = '\n\n'.join(data)
    return content

def extract_pptx(path):
    prs = Presentation(path)
    text = ""
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + "\n"
    return text
